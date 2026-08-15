# -*- coding: utf-8 -*-
"""生成夏令营答辩 PPT。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "paper" / "基于YOLOv12的无人机俯视交通流感知与道路拥堵智能研判系统.pptx"

UI_IMG = ROOT / "paper" / "figures" / "system_ui.png"
CM_IMG = ROOT / "my_uav_vehicle_single" / "my_uav_vehicle_single" / "confusion_matrix.png"
TRAIN_IMG = ROOT / "my_uav_vehicle_single" / "my_uav_vehicle_single" / "results.png"

# 配色
C_PRIMARY = RGBColor(0, 72, 130)
C_ACCENT = RGBColor(0, 160, 210)
C_DARK = RGBColor(30, 40, 55)
C_TEXT = RGBColor(45, 55, 72)
C_LIGHT = RGBColor(245, 248, 252)
C_WHITE = RGBColor(255, 255, 255)
C_MUTED = RGBColor(100, 116, 139)


def set_run(run, *, size: float, bold: bool = False, color: RGBColor = C_TEXT, name: str = "微软雅黑") -> None:
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, prs: Presentation, title: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_PRIMARY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    set_run(run, size=26, bold=True, color=C_WHITE)


def add_bullets(
    slide,
    items: list[str],
    *,
    left: float = 0.7,
    top: float = 1.45,
    width: float = 11.8,
    height: float = 5.5,
    size: float = 20,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.level = 0
        run = p.add_run()
        run.text = item
        set_run(run, size=size, color=C_TEXT)


def add_two_column_bullets(
    slide,
    left_items: list[str],
    right_items: list[str],
    *,
    left_title: str = "",
    right_title: str = "",
) -> None:
    if left_title:
        lt = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(5.5), Inches(0.5))
        r = lt.text_frame.paragraphs[0].add_run()
        r.text = left_title
        set_run(r, size=22, bold=True, color=C_PRIMARY)
    if right_title:
        rt = slide.shapes.add_textbox(Inches(6.8), Inches(1.35), Inches(5.5), Inches(0.5))
        r = rt.text_frame.paragraphs[0].add_run()
        r.text = right_title
        set_run(r, size=22, bold=True, color=C_PRIMARY)
    add_bullets(slide, left_items, left=0.7, top=1.85, width=5.5, height=4.8, size=18)
    add_bullets(slide, right_items, left=6.8, top=1.85, width=5.5, height=4.8, size=18)


def add_image_safe(slide, path: Path, left, top, width) -> bool:
    if path.is_file():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))
        return True
    return False


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_DARK)
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(3.05), prs.slide_width, Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.6))
    p = title_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "基于YOLOv12的无人机俯视交通流感知\n与道路拥堵智能研判系统"
    set_run(run, size=34, bold=True, color=C_WHITE)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(11.5), Inches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Intelligent Road Congestion Assessment System for UAV Top-Down Traffic Flow Perception"
    set_run(r2, size=14, color=C_ACCENT)

    info = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.2))
    tf = info.text_frame
    for i, line in enumerate(
        ["汇报人：（作者姓名）", "单位：（单位名称）", "成果类型：夏令营项目展示 · 2026"]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        set_run(r, size=16, color=C_MUTED)


def slide_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "汇报提纲")
    add_bullets(
        slide,
        [
            "01  研究背景与意义",
            "02  系统总体设计",
            "03  YOLOv12 车辆检测与跟踪",
            "04  道路拥堵指数算法",
            "05  Web 系统实现与展示",
            "06  实验结果与分析",
            "07  总结与展望",
        ],
        size=24,
    )


def slide_background(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "01  研究背景与意义")
    add_two_column_bullets(
        slide,
        [
            "城市化加速，交通拥堵制约出行效率",
            "传统方法依赖路侧检测器 / 浮动车",
            "布设成本高，难以快速机动部署",
            "无人机俯视：80 m、-90° 垂直俯拍",
            "可快速覆盖大范围路网，机动灵活",
        ],
        [
            "实时研判道路运行状态",
            "区分畅通 / 排队 / 走走停停 / 路口等灯",
            "支撑交通规划、信号控制与出行诱导",
            "面向夏令营：可演示、可调试、可导出",
        ],
        left_title="背景",
        right_title="目标",
    )


def slide_challenge(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "问题与挑战")
    add_bullets(
        slide,
        [
            "目标尺度小：UAV 俯视下车辆仅占少量像素",
            "检测框抖动大：逐帧 bbox 波动影响速度/车距估计",
            "路口多队列并存：跨队列大间距会虚增 σdist",
            "需实时运行：消费级 GPU + Web 端同步展示",
            "需可解释：答辩时能说明「Ct 为何升高/降低」",
        ],
        size=22,
    )


def slide_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "02  系统总体设计")
    layers = [
        ("应用层", "FastAPI · MJPEG · WebSocket · 90s 曲线 · CSV 导出"),
        ("算法层", "CongestionCalculator · 车队切分 · 连续状态机 · 抗抖动"),
        ("数据层", "DataCleaner · EMA 平滑 · 3s 滑动窗口 · 30 FPS 时基"),
        ("感知层", "YOLOv12 检测 · ByteTrack 跟踪 · ONNX + DirectML"),
    ]
    y = 1.5
    for name, desc in layers:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(1.2),
            Inches(y),
            Inches(10.8),
            Inches(0.95),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = C_WHITE
        box.line.color.rgb = C_ACCENT
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = f"{name}    "
        set_run(r1, size=20, bold=True, color=C_PRIMARY)
        r2 = p.add_run()
        r2.text = desc
        set_run(r2, size=18, color=C_TEXT)
        y += 1.15


def slide_pipeline(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "数据流程")
    flow = (
        "UAV 俯视视频\n"
        "    ↓\n"
        "YOLOv12 检测 + ByteTrack 跟踪（带 track_id）\n"
        "    ↓\n"
        "DataCleaner：EMA 平滑 · 死区过滤 · 3s 窗口\n"
        "    ↓\n"
        "CongestionCalculator：Vavg · σdist · Ct · 场景标签\n"
        "    ↓\n"
        "Web 实时展示 + 25 列 CSV 导出"
    )
    box = slide.shapes.add_textbox(Inches(1.5), Inches(1.6), Inches(10), Inches(4.8))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = flow
    set_run(run, size=22, color=C_TEXT)


def slide_yolo(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "03  YOLOv12 两阶段迁移学习")
    add_bullets(
        slide,
        [
            "标注：X-AnyLabeling 矩形框标注自采集 UAV 视频帧",
            "阶段一：VisDrone 预训练权重 → 域适应（1280，50 epoch）",
            "阶段二：阶段一 best → 自标注 merged_vehicle 微调（40 epoch）",
            "策略：通用 UAV 检测能力 + 项目场景定制",
            "部署：导出 ONNX · 隔帧推理（每 2 帧 1 次）· DirectML 加速",
        ],
        left=0.7,
        top=1.4,
        width=6.2,
        size=18,
    )
    add_image_safe(slide, TRAIN_IMG, 7.2, 1.35, 5.3)


def slide_metrics(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "检测性能（验证集）")

    metrics = [
        ("VisDrone 域适应", "mAP@0.5  91.2%", "P 94.5%  ·  R 86.3%"),
        ("X-AnyLabeling 微调", "mAP@0.5  97.7%", "P 95.3%  ·  R 95.9%"),
    ]
    x = 0.9
    for title, main, sub in metrics:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.5),
            Inches(5.5),
            Inches(1.5),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = C_PRIMARY
        tf = card.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r0 = p.add_run()
        r0.text = title + "\n"
        set_run(r0, size=16, bold=True, color=C_PRIMARY)
        r1 = p.add_run()
        r1.text = main + "\n"
        set_run(r1, size=28, bold=True, color=C_ACCENT)
        r2 = p.add_run()
        r2.text = sub
        set_run(r2, size=14, color=C_MUTED)
        x += 6.0

    add_image_safe(slide, CM_IMG, 2.0, 3.3, 8.8)
    cap = slide.shapes.add_textbox(Inches(2), Inches(6.85), Inches(9), Inches(0.4))
    r = cap.text_frame.paragraphs[0].add_run()
    r.text = "图：YOLOv12 微调后验证集混淆矩阵"
    set_run(r, size=12, color=C_MUTED)


def slide_algorithm(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "04  道路拥堵指数算法")
    add_bullets(
        slide,
        [
            "特征：Vavg（平均速度）· Dt（帧车距）· σdist（车距波动）· rstop",
            "车队切分：GapThr = k·(Li+Li+1)/2，k=1.5，仅在队内统计车距",
            "指数合成：Ct = 0.67·Sspeed + 0.33·Sfluct（0—100 分）",
            "判定：Ct < 60 畅通  ·  Ct ≥ 60 拥堵",
            "抗抖动：Hold 挂起 · Dist Warmup · 时间基 EMA · 速率限幅 · Episode Reset",
        ],
        size=20,
    )


def slide_scenarios(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "典型场景判别")

    headers = ["场景", "速度特征", "波动特征", "Ct 区间", "标签"]
    rows = [
        ["A 畅通", "Vavg≥Vhigh", "σdist 低", "<60", "畅通"],
        ["B 有序等待", "低速", "σdist≤σlow", "20—40", "有序等待"],
        ["C 走走停停", "低速", "σdist≥σhigh", "≥60", "拥堵"],
        ["E 路口等灯", "静动混合", "静组波动小", "≤30", "路口等灯"],
    ]
    cols, row_n = len(headers), len(rows) + 1
    table = slide.shapes.add_table(row_n, cols, Inches(0.8), Inches(1.55), Inches(11.5), Inches(2.8)).table

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_run(run, size=14, bold=True, color=C_WHITE)

    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run(run, size=13, color=C_TEXT)

    note = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.2))
    tf = note.text_frame
    for i, line in enumerate(
        [
            "实测示例（test.MP4）：17 辆车 · Ct=38.5 · σdist=1.75 → 判定畅通",
            "B/C 区间连续映射，σdist 从 σlow→σhigh 对应 20→80 分，避免阈值跳崖",
            "路口场景 Ct 封顶约 30 分，降低混行误判",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = "• " + line
        set_run(r, size=18, color=C_TEXT)


def slide_web(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "05  Web 系统展示")
    add_bullets(
        slide,
        [
            "FastAPI 后端 · 生产者-消费者多线程 · 帧队列深度 2",
            "MJPEG 视频流 + WebSocket 实时指标推送",
            "四指标卡片：Ct · Vavg · σdist · 车辆数",
            "90 s 滚动曲线 · 60 分阈值虚线 · 参数热更新",
            "逐帧调试模式 + 25 列 CSV 全量导出",
        ],
        left=0.7,
        top=1.35,
        width=4.8,
        height=5.5,
        size=16,
    )
    add_image_safe(slide, UI_IMG, 5.6, 1.25, 6.8)


def slide_experiment(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "06  实验结果与分析")
    add_two_column_bullets(
        slide,
        [
            "平台：Windows · AMD GPU · DirectML",
            "视频：自采集 UAV 俯视场景",
            "场景：直行 / 红灯排队 / 走走停停 / 路口混行",
            "微调 mAP@0.5：97.7%",
        ],
        [
            "Web 端 FPS ≈ 17（隔帧推理）",
            "WebSocket 延迟 < 帧间隔",
            "Ct 可区分四典型场景",
            "CSV 含 raw 分、Dt、Platoons 等 25 列",
        ],
        left_title="实验环境",
        right_title="运行效果",
    )


def slide_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_LIGHT)
    add_top_bar(slide, prs, "07  总结与展望")
    add_bullets(
        slide,
        [
            "完成 YOLOv12 两阶段检测 + ByteTrack + 拥堵指数 + Web 系统集成",
            "贡献：分层迁移训练 · 可解释 Ct 模型 · 可演示工程平台",
            "局限：像素速度（未三维标定）· 单路分析为主",
            "展望：MOTA 跟踪指标 · 多路融合 · 与宏观指数融合",
        ],
        size=22,
    )


def slide_thanks(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, C_DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(1.5))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "谢谢聆听  ·  欢迎提问"
    set_run(run, size=40, bold=True, color=C_WHITE)
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6))
    p2 = sub.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "基于YOLOv12的无人机俯视交通流感知与道路拥堵智能研判系统"
    set_run(r2, size=16, color=C_ACCENT)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_agenda(prs)
    slide_background(prs)
    slide_challenge(prs)
    slide_architecture(prs)
    slide_pipeline(prs)
    slide_yolo(prs)
    slide_metrics(prs)
    slide_algorithm(prs)
    slide_scenarios(prs)
    slide_web(prs)
    slide_experiment(prs)
    slide_summary(prs)
    slide_thanks(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"已生成 PPT: {OUT}")
    print(f"共 {len(prs.slides)} 页")
    return OUT


if __name__ == "__main__":
    build()
